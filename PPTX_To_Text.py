
from nltk.tokenize import  word_tokenize
import os
import pandas as pd
from pathlib import Path
import Library_OF_Functions
from pptx import Presentation

def Convert_PPTX_To_Text(file_name, file_extension, Dir_name, Text_Formats_Dir_name, Datasets_Feautres_Dir_name, Patterns_Data, Data_Summary):
    File_path_For_Patterns_Data=Path(Patterns_Data, "PPTX_Patterns").with_suffix('.csv')
    if(os.path.isfile(File_path_For_Patterns_Data)):
        File_path_For_Risk_Types=Path(Patterns_Data, "Risk_Types").with_suffix('.csv')
        if(os.path.isfile(File_path_For_Risk_Types)):
            open(File_path_For_Risk_Types, 'r')
            C_Dataframe=pd.read_csv(File_path_For_Risk_Types, sep=',')
            Classes_Dataframe = C_Dataframe['Risk_type'].tolist()
            Patterns_Dataframe=pd.read_csv(File_path_For_Patterns_Data, sep=',')
            File_path=Path(Dir_name, file_name).with_suffix(file_extension)
            Presentation_File = Presentation(File_path)
            text=''
            Pattern=[]
            for slide in Presentation_File.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text=text+str(shape.text)
            text = text.encode('utf-8').decode('ascii', 'ignore')
            Words_Tokens=word_tokenize(text)
            Words_Tokens=Library_OF_Functions.Convert_To_Lower(Words_Tokens)
            Words_Tokens=Library_OF_Functions.Stop_removal(Words_Tokens)
            Words_Tokens=Library_OF_Functions.lemmatizer_Process(Words_Tokens)
            if(len(Words_Tokens)>0):
                x,y=Patterns_Dataframe.shape
                for i in range(x):
                    Pattern=Library_OF_Functions.A_Row_Of_Dataframe_to_List(Patterns_Dataframe, i)
                    Percentage=Pattern[len(Pattern)-1]
                    Pattern=Pattern[0:len(Pattern)-1]
                    Pattern=Library_OF_Functions.Convert_To_Lower(Pattern)
                    Pattern_Str = str(Pattern)
                    indexes = [w for w, val in enumerate(Words_Tokens) if val in Pattern_Str]                    
                    List_Of_Words_In_Order, index_Words_In_Order=Library_OF_Functions.Identify_Words_And_Indexes(indexes, Words_Tokens)
                    if len(index_Words_In_Order) >= len(Pattern): 
                        Library_OF_Functions.Find_Patterns_Index_In_Sequence_Order(List_Of_Words_In_Order,index_Words_In_Order , Words_Tokens, Pattern, Datasets_Feautres_Dir_name, file_name, Text_Formats_Dir_name, i+1, Data_Summary, "PPTX", Classes_Dataframe, Patterns_Data)
        else:
            print ("The file risk types is not exist - File Risk_Types.csv is not found in the ", File_path_For_Patterns_Data)
    else:
        print ("The file Patterns of PPTX files is not exist - File PPTX_Patterns.csv is not found in the ", File_path_For_Patterns_Data)